"""KDM6 자동미분 발표자료 (5분, 7장) — 파스텔톤."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── pastel palette ──────────────────────────────────────────────────────────
BG          = RGBColor(0xFA, 0xFA, 0xF7)  # off-white
TITLE_BLUE  = RGBColor(0x35, 0x4B, 0x6F)  # deep slate blue
ACCENT_BLUE = RGBColor(0xA7, 0xC7, 0xE7)  # pastel blue
ACCENT_PINK = RGBColor(0xF8, 0xC7, 0xC9)  # pastel pink
ACCENT_GRN  = RGBColor(0xC7, 0xE5, 0xB4)  # pastel green
ACCENT_YEL  = RGBColor(0xFC, 0xE5, 0xA8)  # pastel yellow
ACCENT_LAV  = RGBColor(0xE0, 0xBB, 0xE4)  # pastel lavender
ACCENT_PEACH= RGBColor(0xFF, 0xD8, 0xB1)  # pastel peach
TEXT_DARK   = RGBColor(0x2C, 0x3E, 0x50)
TEXT_GREY   = RGBColor(0x55, 0x60, 0x70)
LINE_LIGHT  = RGBColor(0xC0, 0xC8, 0xD0)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]


def add_blank_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Pretendard"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_box(slide, x, y, w, h, *, fill=ACCENT_BLUE, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_arrow(slide, x1, y1, x2, y2, *, color=TEXT_GREY, width=2.0):
    cn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    # add arrowhead
    line_xml = cn.line._get_or_add_ln()
    tail = etree.SubElement(line_xml, qn('a:tailEnd'),
                            {'type': 'triangle', 'w': 'med', 'len': 'med'})
    return cn


def add_header(slide, title, subtitle=None):
    add_box(slide, Inches(0), Inches(0), W, Inches(0.95), fill=TITLE_BLUE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.4),
             title, size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.32),
                 subtitle, size=12, color=RGBColor(0xD8, 0xE2, 0xF0),
                 anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_num, total=8):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "KDM6 자동미분 — 새로운 단계의 수치예보로", size=9, color=TEXT_GREY)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=9, color=TEXT_GREY, align=PP_ALIGN.RIGHT)


# ============================================================================
# Slide 1: Title
# ============================================================================
s = add_blank_slide()

# Decorative pastel ribbons
add_box(s, Inches(0), Inches(0), W, Inches(2.0), fill=ACCENT_BLUE)
add_box(s, Inches(0), Inches(5.3), W, Inches(2.2), fill=ACCENT_LAV)

# Title
add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.7),
         "KDM6 자동미분 — 새로운 단계의 수치예보로의 도약",
         size=36, bold=True, color=TITLE_BLUE)

add_text(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.5),
         "Fortran microphysics → PyTorch differentiable oracle",
         size=18, color=TEXT_GREY)

add_text(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.5),
         "물리 과정 자체의 개선과 자료동화 도약을 위한 핵심 기반 기술",
         size=16, color=TEXT_DARK)

# Three pastel callouts at bottom
cw = Inches(3.6); ch = Inches(1.3); cy = Inches(5.6); cx0 = Inches(0.7)
gap = Inches(0.4)

callouts = [
    ("포팅 진행률", "8 / 10", "Step C sub-step 완료", ACCENT_PEACH),
    ("Python 회귀", "111 / 111", "PASS · mac standalone", ACCENT_GRN),
    ("C++ ctest", "6 / 6", "PASS · 2.34s · libtorch", ACCENT_YEL),
]
for i, (label, big, sub, color) in enumerate(callouts):
    bx = cx0 + i * (cw + gap)
    add_box(s, bx, cy, cw, ch, fill=color)
    add_text(s, bx, cy + Inches(0.1), cw, Inches(0.3), label,
             size=11, color=TEXT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, bx, cy + Inches(0.4), cw, Inches(0.55), big,
             size=28, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, bx, cy + Inches(0.95), cw, Inches(0.3), sub,
             size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4),
         "2026-04-28 · KDM6_AD project (KIM-meso v1.0 slot 47)",
         size=10, color=TEXT_GREY, align=PP_ALIGN.LEFT)


# ============================================================================
# Slide 2: 배경 + 전/후 비교
# ============================================================================
s = add_blank_slide()
add_header(s, "왜 KDM6를 자동미분 가능하게 만드는가",
           "Fortran 단방향 forward → PyTorch 양방향 (forward + 미분)")

# Left: 기존 Fortran (BEFORE)
left_x = Inches(0.5); left_y = Inches(1.3); left_w = Inches(6.0); left_h = Inches(5.4)
add_box(s, left_x, left_y, left_w, left_h, fill=ACCENT_PINK)
add_text(s, left_x + Inches(0.3), left_y + Inches(0.2), left_w - Inches(0.6), Inches(0.45),
         "BEFORE — Fortran KDM6 (4,281 라인)", size=18, bold=True, color=TITLE_BLUE)
add_text(s, left_x + Inches(0.3), left_y + Inches(0.7), left_w - Inches(0.6), Inches(0.35),
         "단방향 forward만 가능 · 미분 정보 없음",
         size=12, color=TEXT_GREY)

bullets_before = [
    "if-branch 312개 · 모든 cliff에서 미분 0 또는 정의 안 됨",
    "분모/제곱근 보호 없음 — IEEE invalid 수동 처리 (86건)",
    "TLM/ADJ는 별도 Tapenade 산물 (a_*, pushreal8/popreal8) 적층",
    "매개변수 보정 = 수많은 forward run의 grid search",
    "DA에 사용 시: adjoint를 다시 별도 빌드해야 함",
]
by = left_y + Inches(1.2)
for txt in bullets_before:
    add_text(s, left_x + Inches(0.4), by, left_w - Inches(0.7), Inches(0.4),
             "•  " + txt, size=12, color=TEXT_DARK)
    by += Inches(0.55)

# Right: 새 PyTorch (AFTER)
right_x = Inches(6.83); right_y = Inches(1.3); right_w = Inches(6.0); right_h = Inches(5.4)
add_box(s, right_x, right_y, right_w, right_h, fill=ACCENT_GRN)
add_text(s, right_x + Inches(0.3), right_y + Inches(0.2), right_w - Inches(0.6), Inches(0.45),
         "AFTER — PyTorch / libtorch oracle", size=18, bold=True, color=TITLE_BLUE)
add_text(s, right_x + Inches(0.3), right_y + Inches(0.7), right_w - Inches(0.6), Inches(0.35),
         "forward + VJP/JVP/Jacobian 자동 산출",
         size=12, color=TEXT_GREY)

bullets_after = [
    "모든 분기를 torch.where 마스크 + safe_div/safe_pow로 미분 가능",
    "동일 코드로 forward, adjoint, tangent 즉시 사용 (graph 자동)",
    "매개변수 (PEAUT, NCRK1/2 등) 자동 grad — 보정 비용 ↓↓",
    "DA 결합: 4D-Var adjoint 운영 모델의 1차 입력으로 즉시 사용",
    "Python ↔ C++ 양쪽 동일 oracle (mac/linux 동일 검증)",
]
by = right_y + Inches(1.2)
for txt in bullets_after:
    add_text(s, right_x + Inches(0.4), by, right_w - Inches(0.7), Inches(0.4),
             "•  " + txt, size=12, color=TEXT_DARK)
    by += Inches(0.55)

# Center arrow
add_arrow(s, Inches(6.5), Inches(4.0), Inches(6.85), Inches(4.0), color=TITLE_BLUE, width=4.0)

add_footer(s, 2)


# ============================================================================
# Slide 3: 포팅 과정 (7-step plan, 진행률)
# ============================================================================
s = add_blank_slide()
add_header(s, "포팅 과정 — kdm62D 7-단계 분해",
           "Fortran 본문 ~2,630 라인을 의존성 graph 기반으로 분해 → 단계적 검증")

# Step bar
step_data = [
    ("Step 0", "foundations\n(constants/ops/state)", ACCENT_GRN, "✓"),
    ("Step 0a", "slope\n(rain/snow/graupel/ice)", ACCENT_GRN, "✓"),
    ("Step A", "ProgB_param\n(graupel density 진단)", ACCENT_GRN, "✓"),
    ("Step B", "warm rates + sat adj\n(B1-B5)", ACCENT_GRN, "✓"),
    ("Step C", "ice phase rates\n(C1-C4 ✓ / C5-C6)", ACCENT_YEL, "8/10"),
    ("Step D", "melting / freezing\n(Bigg, t40c, t0c)", ACCENT_PEACH, "☐"),
    ("Step E", "sedimentation\n(NISLFV-PLM)", ACCENT_PEACH, "☐"),
    ("Step F", "kdm62D coordinator\n(A~E wire)", ACCENT_LAV, "☐"),
    ("Step G", "KIM-meso integration\n(slot 47, ISO_C_BINDING)", ACCENT_LAV, "☐"),
]
n = len(step_data)
strip_x = Inches(0.5); strip_y = Inches(1.4); strip_w = Inches(12.3); strip_h = Inches(1.6)
gap = Inches(0.05)
box_w = (strip_w - gap * (n - 1)) / n
for i, (step, body, color, mark) in enumerate(step_data):
    bx = strip_x + i * (box_w + gap)
    add_box(s, bx, strip_y, box_w, strip_h, fill=color)
    add_text(s, bx, strip_y + Inches(0.1), box_w, Inches(0.3),
             step, size=11, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, bx, strip_y + Inches(0.4), box_w, Inches(0.7),
             body, size=9, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, bx, strip_y + Inches(1.15), box_w, Inches(0.4),
             mark, size=18, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)

# Cycle box (below)
cy_box_y = Inches(3.3)
add_box(s, Inches(0.5), cy_box_y, Inches(12.3), Inches(1.0), fill=ACCENT_BLUE)
add_text(s, Inches(0.7), cy_box_y + Inches(0.1), Inches(11.9), Inches(0.4),
         "각 단계의 사이클", size=14, bold=True, color=TITLE_BLUE)

cycle_steps = ["Fortran 읽기", "Python oracle", "Python 테스트",
               "C++ libtorch 포트", "C++ ctest", "wiki 적재"]
cx_y = cy_box_y + Inches(0.5)
cx0 = Inches(0.8); cw = Inches(2.0)
for i, txt in enumerate(cycle_steps):
    bx = cx0 + Inches(i * 2.05)
    add_text(s, bx, cx_y, cw, Inches(0.4), txt,
             size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    if i < len(cycle_steps) - 1:
        add_text(s, bx + cw - Inches(0.05), cx_y, Inches(0.15), Inches(0.4),
                 "→", size=14, color=TITLE_BLUE)

# Bottom: 검증 통계
v_y = Inches(4.6)
add_text(s, Inches(0.5), v_y, Inches(12.3), Inches(0.4),
         "현재까지 검증 (mac, libtorch from PyTorch wheel)",
         size=14, bold=True, color=TITLE_BLUE)

stat_y = v_y + Inches(0.6)
stats = [
    ("Python pytest 회귀", "111 / 111 PASS", ACCENT_GRN),
    ("C++ ctest 6 모듈", "6 / 6 PASS · 2.34s", ACCENT_YEL),
    ("발견한 잠재 결함", "muc=2 sigma_c NaN, eac 부호\nC++ in-place AD 충돌 (수정)", ACCENT_PINK),
]
sw = Inches(4.0); sh = Inches(1.6); sx0 = Inches(0.5); sgap = Inches(0.15)
for i, (lbl, big, color) in enumerate(stats):
    bx = sx0 + i * (sw + sgap)
    add_box(s, bx, stat_y, sw, sh, fill=color)
    add_text(s, bx, stat_y + Inches(0.15), sw, Inches(0.35), lbl,
             size=12, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.1), stat_y + Inches(0.55), sw - Inches(0.2), Inches(0.95), big,
             size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 3)


# ============================================================================
# Slide 4 (NEW): 자동미분 · 연산 그래프 개념
# ============================================================================
s = add_blank_slide()
add_header(s, "자동미분의 작동 원리 — 연산 그래프 + JVP / VJP",
           "PyTorch가 forward 시 그래프를 기록 → backward 시 chain rule 자동 적용")

# ── Left: computation graph 도식 ─────────────────────────────────────────
gx = Inches(0.5); gy = Inches(1.3); gw = Inches(7.0); gh = Inches(5.4)
add_box(s, gx, gy, gw, gh, fill=ACCENT_BLUE)
add_text(s, gx + Inches(0.3), gy + Inches(0.2), gw - Inches(0.6), Inches(0.4),
         "연산 그래프 (computation graph)",
         size=18, bold=True, color=TITLE_BLUE)

# Node positions (centers)
node_w = Inches(1.15); node_h = Inches(0.55)
ny_fwd = gy + Inches(1.5)   # forward chain row
ny_bwd = gy + Inches(3.4)   # backward chain row
node_xs = [Inches(1.0), Inches(2.55), Inches(4.1), Inches(5.65)]

# Forward chain: x → f₁ → z₁ → f₂ → z₂ → ... → y
fwd_labels = ["x\n(state)", "z₁", "z₂", "y\n(rate)"]
for i, (nx, lbl) in enumerate(zip(node_xs, fwd_labels)):
    color = ACCENT_GRN if i in (0, len(fwd_labels) - 1) else RGBColor(0xFF, 0xFF, 0xFF)
    add_box(s, nx, ny_fwd, node_w, node_h, fill=color, line=TITLE_BLUE)
    add_text(s, nx, ny_fwd, node_w, node_h, lbl,
             size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Forward arrows + labels (f₁, f₂, ...)
for i in range(len(fwd_labels) - 1):
    x1 = node_xs[i] + node_w
    x2 = node_xs[i + 1]
    y_mid = ny_fwd + node_h / 2
    add_arrow(s, x1, y_mid, x2, y_mid, color=TITLE_BLUE, width=2.0)
    # label above arrow
    add_text(s, x1, ny_fwd - Inches(0.3), x2 - x1, Inches(0.3),
             f"f{['₁','₂','₃'][i]}", size=11, bold=True,
             color=TITLE_BLUE, align=PP_ALIGN.CENTER)

# "forward" label
add_text(s, gx + Inches(0.3), ny_fwd - Inches(0.7), Inches(2.5), Inches(0.3),
         "▶ forward (값 + 그래프 기록)",
         size=11, bold=True, color=TITLE_BLUE)

# Backward chain (same nodes, reversed arrows)
bwd_labels = ["∂L/∂x", "∂L/∂z₁", "∂L/∂z₂", "∂L/∂y"]
for i, (nx, lbl) in enumerate(zip(node_xs, bwd_labels)):
    color = ACCENT_PINK if i in (0, len(bwd_labels) - 1) else RGBColor(0xFF, 0xFF, 0xFF)
    add_box(s, nx, ny_bwd, node_w, node_h, fill=color, line=ACCENT_PINK)
    add_text(s, nx, ny_bwd, node_w, node_h, lbl,
             size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Backward arrows (reversed)
for i in range(len(bwd_labels) - 1, 0, -1):
    x1 = node_xs[i]
    x2 = node_xs[i - 1] + node_w
    y_mid = ny_bwd + node_h / 2
    add_arrow(s, x1, y_mid, x2, y_mid, color=ACCENT_PINK, width=2.0)

# "backward" label
add_text(s, gx + Inches(0.3), ny_bwd - Inches(0.4), Inches(3.0), Inches(0.3),
         "◀ backward (chain rule 자동 적용)",
         size=11, bold=True, color=ACCENT_PINK)

# Caption note (bottom of left box)
note_y = gy + Inches(4.3)
add_text(s, gx + Inches(0.3), note_y, gw - Inches(0.6), Inches(0.35),
         "PyTorch는 forward 시 모든 연산을 노드/엣지로 기록.",
         size=11, color=TEXT_DARK)
add_text(s, gx + Inches(0.3), note_y + Inches(0.4), gw - Inches(0.6), Inches(0.35),
         "loss.backward() 호출 시 그래프를 *역방향* 순회하며 각 입력의 grad 누적.",
         size=11, color=TEXT_DARK)
add_text(s, gx + Inches(0.3), note_y + Inches(0.85), gw - Inches(0.6), Inches(0.35),
         "→ KDM6 oracle은 모든 분기를 미분 가능하게 표현하여 그래프가 끊기지 않도록 함.",
         size=11, bold=True, color=TITLE_BLUE)

# ── Right: JVP & VJP cards ───────────────────────────────────────────────
rx = Inches(7.83); rw = Inches(5.0)

# JVP card
jx_y = Inches(1.3); jx_h = Inches(2.55)
add_box(s, rx, jx_y, rw, jx_h, fill=ACCENT_GRN)
add_text(s, rx + Inches(0.25), jx_y + Inches(0.15), rw - Inches(0.5), Inches(0.4),
         "JVP — Forward-mode (Jacobian × Vector)",
         size=14, bold=True, color=TITLE_BLUE)
add_text(s, rx + Inches(0.25), jx_y + Inches(0.6), rw - Inches(0.5), Inches(0.45),
         "J · v   :   ∂y/∂x · v   →   y의 perturbation",
         size=14, bold=True, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), jx_y + Inches(1.15), rw - Inches(0.5), Inches(0.4),
         "▸ 한 perturbation v → 모든 출력의 변화",
         size=11, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), jx_y + Inches(1.5), rw - Inches(0.5), Inches(0.4),
         "▸ 비용 ≈ forward 1회 (cheaper for many outputs)",
         size=11, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), jx_y + Inches(1.85), rw - Inches(0.5), Inches(0.6),
         "▸ KDM6 활용: EnKF perturbation propagation,\n   tangent linear model (TLM)",
         size=11, bold=True, color=TITLE_BLUE)

# VJP card
vx_y = Inches(4.1); vx_h = Inches(2.6)
add_box(s, rx, vx_y, rw, vx_h, fill=ACCENT_PINK)
add_text(s, rx + Inches(0.25), vx_y + Inches(0.15), rw - Inches(0.5), Inches(0.4),
         "VJP — Reverse-mode (Vector × Jacobian)",
         size=14, bold=True, color=TITLE_BLUE)
add_text(s, rx + Inches(0.25), vx_y + Inches(0.6), rw - Inches(0.5), Inches(0.45),
         "uᵀ · J   :   uᵀ · ∂y/∂x   →   x의 grad",
         size=14, bold=True, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), vx_y + Inches(1.15), rw - Inches(0.5), Inches(0.4),
         "▸ 한 scalar loss → 모든 입력의 grad 한 번에",
         size=11, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), vx_y + Inches(1.5), rw - Inches(0.5), Inches(0.4),
         "▸ 비용 ≈ forward 1회 (cheaper for many inputs)",
         size=11, color=TEXT_DARK)
add_text(s, rx + Inches(0.25), vx_y + Inches(1.85), rw - Inches(0.5), Inches(0.6),
         "▸ KDM6 활용: 4D-Var adjoint, 매개변수 보정,\n   관측-모델 잔차 → 상태/매개변수 grad",
         size=11, bold=True, color=TITLE_BLUE)

# Bottom note (full width)
add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.3),
         "Jacobian 전체 행렬은 비싸지만 (output × input), JVP/VJP는 *벡터 한 번*만 필요 — DA·보정에 핵심",
         size=11, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)

add_footer(s, 4)


# ============================================================================
# Slide 5: 자동미분 핵심 조치
# ============================================================================
s = add_blank_slide()
add_header(s, "자동미분을 위한 핵심 조치 · 주의 사항",
           "Fortran 직역만으로는 부족 — graph 보존을 위한 idiom 필수")

# Two columns
left_x = Inches(0.5); col_w = Inches(6.15); col_h = Inches(5.4); top_y = Inches(1.3)

# Left: 핵심 idiom
add_box(s, left_x, top_y, col_w, col_h, fill=ACCENT_BLUE)
add_text(s, left_x + Inches(0.3), top_y + Inches(0.2), col_w - Inches(0.6), Inches(0.4),
         "필수 idiom", size=18, bold=True, color=TITLE_BLUE)

idioms = [
    ("`.item()` 0회",
     "tensor → Python scalar 변환은 graph 단절. 부득이 시 `with torch.no_grad():`"),
    ("`torch.where` 마스크",
     "if-else를 미분 가능 분기로. inactive cell의 graph도 깔끔히 차단"),
    ("safe_div_pos · safe_pow",
     "분모 `clamp(min=EPS)`, 음수 base/exponent 보호. forward/backward 둘 다 안전"),
    ("`abs(vt2 - vt2)` 류",
     "subgradient cliff at equality — 일반 관측 X. 매끄러움 필요시 sigmoid blend"),
    ("complete sublim/evap mutation",
     "Fortran in-place mutation은 oracle에서 `*_adj` 별도 출력으로 풀어내기"),
    ("C++ in-place op 금지",
     "`logical_and_`, `logical_or_` → `WhereBackward` version mismatch. out-of-place 사용"),
]
iy = top_y + Inches(0.7)
for ttl, body in idioms:
    add_text(s, left_x + Inches(0.4), iy, col_w - Inches(0.7), Inches(0.3),
             "▶  " + ttl, size=12, bold=True, color=TITLE_BLUE)
    add_text(s, left_x + Inches(0.7), iy + Inches(0.32), col_w - Inches(1.0), Inches(0.4),
             body, size=10, color=TEXT_DARK)
    iy += Inches(0.78)

# Right: 발견 사례
right_x = Inches(6.83)
add_box(s, right_x, top_y, col_w, col_h, fill=ACCENT_PEACH)
add_text(s, right_x + Inches(0.3), top_y + Inches(0.2), col_w - Inches(0.6), Inches(0.4),
         "포팅 중 발견한 잠재 결함", size=18, bold=True, color=TITLE_BLUE)

discoveries = [
    ("muc=2에서 g6pmc - g3pmc² = -0.5",
     "Fortran (-0.5)^(1/6) → silent NaN. IEEE halting 끔 + finite-repair 없음"),
    ("eac 부호 (직관 반대)",
     "exp(0.07·(-supcol)): cold일수록 collection eff *감소*. 운영 의도 검증 영역"),
    ("nracs commented-out",
     "Fortran 2125-2136 주석 처리. 과거 활성 → 현재 비활성의 코드 archaeology"),
    ("paacw double-decrement",
     "HM에서 snow → graupel sequential mass transfer. caller가 mass-balance 명시"),
    ("psacr supcol gate 누락",
     "psacr는 cold/warm 모두 active. 일반 직관(snow는 cold 전용)과 다름"),
]
iy = top_y + Inches(0.7)
for ttl, body in discoveries:
    add_text(s, right_x + Inches(0.4), iy, col_w - Inches(0.7), Inches(0.3),
             "•  " + ttl, size=12, bold=True, color=TITLE_BLUE)
    add_text(s, right_x + Inches(0.7), iy + Inches(0.32), col_w - Inches(1.0), Inches(0.4),
             body, size=10, color=TEXT_DARK)
    iy += Inches(0.85)

# Bottom note
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
         "Oracle 직역 자체가 *코드 검증 도구*가 됨 — 미분 안전성 보장 + Fortran 운영 의도 자동 surface",
         size=11, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)

add_footer(s, 5)


# ============================================================================
# Slide 6: 활용 1 — 물리 과정 자체의 개선
# ============================================================================
s = add_blank_slide()
add_header(s, "활용 ① — 물리 과정 자체의 개선",
           "자동미분으로 *기존 KDM6의 산식·매개변수 자체를 데이터로 보정*")

# Three pastel cards
cards = [
    ("매개변수 자동 보정",
     "PEAUT, NCRK1/2, ECCBRK 등\n이전: 수많은 forward run의 grid search\n이후: ∂Loss/∂param 직접 계산\n관측-모델 잔차로 *최적값* 자동 도출",
     ACCENT_GRN),
    ("산식 의도 검증",
     "포팅 중 발견한 muc=2 NaN, eac 부호 등\n자동미분 oracle은 산식 자체의\n*수치 안정성 + 의도-구현 일치*를\n자동으로 surface",
     ACCENT_LAV),
    ("분기 매끄러움 결정",
     "물리 분기(big/small drop, T<T0 등)\n→ sigmoid blend 도입 시 미분 신뢰성 ↑\n수치 가드는 stop-gradient 유지\n둘의 분리가 *분기 의미*를 명확히 함",
     ACCENT_BLUE),
]

card_w = Inches(4.0); card_h = Inches(4.5); cy = Inches(1.4); cx0 = Inches(0.5); gap = Inches(0.15)
for i, (ttl, body, color) in enumerate(cards):
    bx = cx0 + i * (card_w + gap)
    add_box(s, bx, cy, card_w, card_h, fill=color)
    add_text(s, bx + Inches(0.3), cy + Inches(0.3), card_w - Inches(0.6), Inches(0.6),
             ttl, size=18, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    # divider line
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              bx + Inches(0.5), cy + Inches(1.0),
                              card_w - Inches(1.0), Emu(7000))
    div.fill.solid()
    div.fill.fore_color.rgb = TITLE_BLUE
    div.line.fill.background()
    add_text(s, bx + Inches(0.3), cy + Inches(1.1), card_w - Inches(0.6), Inches(3.2),
             body, size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT)

# Bottom flow
fy = Inches(6.05)
add_box(s, Inches(0.5), fy, Inches(12.3), Inches(0.85), fill=ACCENT_YEL)
add_text(s, Inches(0.6), fy + Inches(0.1), Inches(12.1), Inches(0.32),
         "보정 사이클 — *연 단위 → 일 단위로 단축* 가능",
         size=13, bold=True, color=TITLE_BLUE)
add_text(s, Inches(0.6), fy + Inches(0.42), Inches(12.1), Inches(0.4),
         "관측 자료 → forward run → ∂Loss/∂param (자동) → 매개변수 갱신 → 재실험",
         size=12, color=TEXT_DARK)

add_footer(s, 6)


# ============================================================================
# Slide 7: 활용 2 — 자료동화 개선
# ============================================================================
s = add_blank_slide()
add_header(s, "활용 ② — 자료동화 (DA) 개선",
           "KIM-meso slot 47 — 운영 forward(slot 37)와 *나란히 구동*")

# Architecture diagram
diag_y = Inches(1.4); diag_h = Inches(3.5)

# 3 vertical lanes
lane_w = Inches(3.9); lane_x = [Inches(0.5), Inches(4.6), Inches(8.7)]
lane_titles = ["관측 정보", "KIM-meso 운영", "DA · 분석"]
lane_colors = [ACCENT_PEACH, ACCENT_BLUE, ACCENT_LAV]

for i, (x, ttl, color) in enumerate(zip(lane_x, lane_titles, lane_colors)):
    add_box(s, x, diag_y, lane_w, diag_h, fill=color)
    add_text(s, x, diag_y + Inches(0.15), lane_w, Inches(0.4),
             ttl, size=14, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)

# Inside boxes
def lane_box(slide, lane_idx, sub_y, h, label, color=RGBColor(0xFF, 0xFF, 0xFF)):
    x = lane_x[lane_idx] + Inches(0.25)
    w = lane_w - Inches(0.5)
    add_box(slide, x, sub_y, w, h, fill=color, line=TEXT_GREY)
    add_text(slide, x, sub_y, w, h, label, size=11, color=TEXT_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x, w

# 관측 boxes
ox, ow = lane_box(s, 0, diag_y + Inches(0.7), Inches(0.6), "위성 (radiance)")
lane_box(s, 0, diag_y + Inches(1.5), Inches(0.6), "지상·항공 관측")
lane_box(s, 0, diag_y + Inches(2.3), Inches(0.6), "레이더 반사도")

# 운영 boxes
lane_box(s, 1, diag_y + Inches(0.7), Inches(0.6), "Slot 37 — Fortran KDM6\n(비트-재현 forward)",
         color=RGBColor(0xFF, 0xFF, 0xFF))
lane_box(s, 1, diag_y + Inches(1.6), Inches(0.6), "Slot 47 — PyTorch KDM6\n(미분가능 oracle)", color=ACCENT_GRN)
lane_box(s, 1, diag_y + Inches(2.5), Inches(0.6), "side-by-side 출력 비교")

# DA boxes
lane_box(s, 2, diag_y + Inches(0.7), Inches(0.6), "4D-Var · adjoint")
lane_box(s, 2, diag_y + Inches(1.5), Inches(0.6), "EnKF · tangent linear")
lane_box(s, 2, diag_y + Inches(2.3), Inches(0.6), "매개변수 동시 추정")

# Arrows: 관측 → slot 47, slot 47 → DA
add_arrow(s, lane_x[0] + lane_w, diag_y + Inches(1.0),
          lane_x[1], diag_y + Inches(1.9), color=TITLE_BLUE, width=2.5)
add_arrow(s, lane_x[1] + lane_w, diag_y + Inches(1.9),
          lane_x[2], diag_y + Inches(1.0), color=TITLE_BLUE, width=2.5)

# Bottom benefits
b_y = Inches(5.2)
benefits = [
    ("Adjoint 운영 비용 ↓",
     "별도 Tapenade 빌드 불필요\n같은 oracle에서 VJP 즉시"),
    ("매개변수 + 상태 동시 추정",
     "기존 DA: 상태만 분석\n→ 매개변수까지 *joint* 추정"),
    ("관측-모델 미스매치 진단",
     "Jacobian 직접 추출 → 어느 process가\n특정 관측에 가장 민감한지 추적"),
]
bw = Inches(4.0); bh = Inches(1.7); by = b_y; bx0 = Inches(0.5); bgap = Inches(0.15)
for i, (ttl, body) in enumerate(benefits):
    bx = bx0 + i * (bw + bgap)
    add_box(s, bx, by, bw, bh, fill=ACCENT_YEL)
    add_text(s, bx + Inches(0.2), by + Inches(0.15), bw - Inches(0.4), Inches(0.4),
             ttl, size=13, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.2), by + Inches(0.6), bw - Inches(0.4), Inches(1.0),
             body, size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_footer(s, 7)


# ============================================================================
# Slide 8: 비전 + 향후 계획
# ============================================================================
s = add_blank_slide()
add_header(s, "비전 — 새로운 단계의 수치예보로",
           "KDM6 자동미분이 가져올 운영 / 연구 / 도약의 세 축")

# Three vertical sectors with icons (text-based)
sw = Inches(4.0); sh = Inches(3.4); sy = Inches(1.4); sx0 = Inches(0.5); sgap = Inches(0.15)

sectors = [
    ("운영", "OPERATIONAL",
     "▸ DA 시스템에 미분 정보 직결\n▸ 매개변수 보정 사이클 단축\n▸ slot 37 비트-재현성 보존\n▸ Python ↔ C++ 동시 운영",
     ACCENT_GRN),
    ("연구", "RESEARCH",
     "▸ 산식 의도-구현 mismatch 발굴\n▸ 분기 의미(물리 vs 수치) 명확화\n▸ 새 매개변수화 도입 비용 ↓\n▸ ML hybrid (NN closure) 연결 가능",
     ACCENT_LAV),
    ("도약", "LEAP",
     "▸ Differentiable atmospheric model\n▸ End-to-end 관측-모델 학습\n▸ 4D-Var 확장 + 매개변수 추정\n▸ KIM 수치예보의 차세대 entry",
     ACCENT_PEACH),
]
for i, (ttl, sub, body, color) in enumerate(sectors):
    bx = sx0 + i * (sw + sgap)
    add_box(s, bx, sy, sw, sh, fill=color)
    add_text(s, bx + Inches(0.3), sy + Inches(0.25), sw - Inches(0.6), Inches(0.6),
             ttl, size=24, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.3), sy + Inches(0.85), sw - Inches(0.6), Inches(0.3),
             sub, size=10, color=TEXT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.3), sy + Inches(1.3), sw - Inches(0.6), Inches(2.0),
             body, size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT)

# Bottom: 향후 계획
fy = Inches(5.05)
add_box(s, Inches(0.5), fy, Inches(12.3), Inches(1.9), fill=ACCENT_BLUE)
add_text(s, Inches(0.7), fy + Inches(0.15), Inches(11.9), Inches(0.4),
         "향후 계획 (단기 → 중기)",
         size=14, bold=True, color=TITLE_BLUE)

# Timeline arrow
ty = fy + Inches(0.7)
arrow_x1 = Inches(0.9); arrow_x2 = Inches(12.4)
add_arrow(s, arrow_x1, ty + Inches(0.5), arrow_x2, ty + Inches(0.5),
          color=TITLE_BLUE, width=4.0)

steps_plan = [
    ("Step C 마무리", "C5 + C6"),
    ("Step D · E", "melting/freezing,\nsedimentation"),
    ("Step F coordinator", "kdm62D wire"),
    ("Step G", "KIM-meso 통합"),
    ("DA 결합 실험", "관측 사례 검증"),
]
n = len(steps_plan)
sx0 = Inches(0.9); spacing = (arrow_x2 - arrow_x1) / (n - 1)
for i, (ttl, sub) in enumerate(steps_plan):
    cx = sx0 + i * spacing
    # dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.1), ty + Inches(0.4),
                              Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = TITLE_BLUE
    dot.line.fill.background()
    add_text(s, cx - Inches(1.1), ty - Inches(0.05), Inches(2.2), Inches(0.4),
             ttl, size=11, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, cx - Inches(1.1), ty + Inches(0.75), Inches(2.2), Inches(0.5),
             sub, size=9, color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_footer(s, 8)


# Save
out_path = "/Users/yhlee/KDM6AD/KDM6_AD_concept_v2.pptx"
prs.save(out_path)
print(f"saved: {out_path}")
